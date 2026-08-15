"""
File Agent — PDF reading, DOCX/PPTX generation, report creation.
"""
from __future__ import annotations
from pathlib import Path
from typing import Any, Dict
import structlog
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from app.agents.base import BaseAgent
from app.core.langgraph.state import AgentState

log = structlog.get_logger(__name__)


class FileAgent(BaseAgent):
    name = "file_agent"
    description = "Read PDFs, generate reports, create DOCX and PPTX"
    supported_intents = ["file_read", "file_write", "pdf_analysis"]

    async def execute(self, state: AgentState) -> Dict[str, Any]:
        action = state.metadata.get("file_action", "read_pdf")

        if action == "read_pdf":
            return await self._read_pdf(state)
        elif action == "analyze_pdf":
            return await self._analyze_pdf(state)
        elif action == "create_docx":
            return await self._create_docx(state)
        elif action == "create_pptx":
            return await self._create_pptx(state)
        elif action == "generate_report":
            return await self._generate_report(state)
        return {"final_response": f"Unknown file action: {action}"}

    async def _read_pdf(self, state: AgentState) -> Dict[str, Any]:
        file_path = state.metadata.get("file_path", "")
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            text = "\n".join(p.extract_text() or "" for p in reader.pages)
            return {
                "final_response": f"📄 **PDF Content** ({len(reader.pages)} pages):\n\n{text[:5000]}{'...' if len(text) > 5000 else ''}",
                "artifacts": [{"type": "pdf_content", "content": text, "pages": len(reader.pages)}],
            }
        except Exception as e:
            return {"error": str(e), "final_response": f"PDF read error: {e}"}

    async def _analyze_pdf(self, state: AgentState) -> Dict[str, Any]:
        read_result = await self._read_pdf(state)
        if read_result.get("error"):
            return read_result
        text = read_result["artifacts"][0]["content"]
        prompt = ChatPromptTemplate.from_messages([
            ("system", "Analyze this document and provide: 1) Summary, 2) Key points, 3) Topics, 4) Important data/figures"),
            ("human", "Document:\n{text}\n\nAnalysis:"),
        ])
        chain = prompt | self.llm | StrOutputParser()
        analysis = await chain.ainvoke({"text": text[:12000]})
        return {"final_response": analysis, "response_type": "report"}

    async def _create_docx(self, state: AgentState) -> Dict[str, Any]:
        try:
            from docx import Document
            from docx.shared import Inches, Pt
            doc = Document()
            title = state.metadata.get("title", "Document")
            content = state.metadata.get("content", state.raw_input)
            doc.add_heading(title, 0)
            for paragraph in content.split("\n\n"):
                if paragraph.strip():
                    doc.add_paragraph(paragraph.strip())
            output_path = state.metadata.get("output", f"{title.replace(' ', '_')}.docx")
            doc.save(output_path)
            return {"final_response": f"✅ DOCX created: `{output_path}`", "response_type": "action"}
        except Exception as e:
            return {"error": str(e), "final_response": f"DOCX error: {e}"}

    async def _create_pptx(self, state: AgentState) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            prs = Presentation()
            slides_data = state.metadata.get("slides", [{"title": state.metadata.get("title", "Presentation"), "content": state.raw_input}])
            for slide_data in slides_data:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = slide_data.get("title", "")
                slide.placeholders[1].text = slide_data.get("content", "")
            output_path = state.metadata.get("output", "presentation.pptx")
            prs.save(output_path)
            return {"final_response": f"✅ PPTX created: `{output_path}` ({len(slides_data)} slides)", "response_type": "action"}
        except Exception as e:
            return {"error": str(e), "final_response": f"PPTX error: {e}"}

    async def _generate_report(self, state: AgentState) -> Dict[str, Any]:
        prompt = ChatPromptTemplate.from_messages([
            ("system", "You are a professional report writer. Generate a comprehensive, well-structured report."),
            ("human", "Report topic/data: {input}\n\nGenerate a professional report with executive summary, findings, and recommendations."),
        ])
        chain = prompt | self.llm | StrOutputParser()
        report = await chain.ainvoke({"input": state.raw_input})
        return {"final_response": report, "response_type": "report"}
